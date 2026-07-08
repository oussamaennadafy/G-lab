#!/usr/bin/env python3
"""Génère la présentation PowerPoint professionnelle du projet G-Lab."""

import sys
from pathlib import Path

try:
    from pptx import Presentation
except ImportError:
    print("Erreur : le module 'python-pptx' n'est pas installé.", file=sys.stderr)
    print("", file=sys.stderr)
    print("Utilisez le script wrapper (recommandé) :", file=sys.stderr)
    print("  bash livrables/generer_presentation.sh", file=sys.stderr)
    print("", file=sys.stderr)
    print("Ou installez manuellement dans un venv :", file=sys.stderr)
    print("  python3 -m venv .venv-pptx", file=sys.stderr)
    print("  .venv-pptx/bin/pip install -r livrables/requirements.txt", file=sys.stderr)
    print("  .venv-pptx/bin/python livrables/generer_presentation.py", file=sys.stderr)
    sys.exit(1)

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

OUTPUT = "livrables/Presentation_G-Lab.pptx"
AUTEUR = "Oussama Ennadafy"
PROJET = "G-Lab"
DATE = "Juillet 2026"

# Palette professionnelle
NAVY       = RGBColor(0x0D, 0x1B, 0x2A)   # Fond sombre
NAVY_MID   = RGBColor(0x1B, 0x3A, 0x5C)   # Bandeau header
ACCENT     = RGBColor(0x00, 0x96, 0xC7)   # Bleu accent (tech)
ACCENT2    = RGBColor(0x48, 0xCA, 0xE4)   # Bleu clair
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT      = RGBColor(0xF0, 0xF4, 0xF8)   # Fond clair slides contenu
TEXT_DARK  = RGBColor(0x1E, 0x29, 0x3B)
TEXT_MID   = RGBColor(0x47, 0x55, 0x69)
TEXT_LIGHT = RGBColor(0x94, 0xA3, 0xB8)
GOLD       = RGBColor(0xF5, 0x9E, 0x0B)   # Accent doré (highlights)

SLIDE_W = Inches(13.333)   # 16:9 widescreen
SLIDE_H = Inches(7.5)
MARGIN  = Inches(0.6)
HEADER_H = Inches(1.05)
FOOTER_H = Inches(0.45)


# ── Helpers graphiques ──────────────────────────────────────────────

def _fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def _no_line(shape):
    shape.line.fill.background()


def _set_font(run, size, color=TEXT_DARK, bold=False, italic=False):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = "Calibri"


def _add_rect(slide, left, top, width, height, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    _fill(s, color)
    _no_line(s)
    return s


def _add_rounded(slide, left, top, width, height, color):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    _fill(s, color)
    _no_line(s)
    return s


def _textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def _write_paragraph(tf, text, size=16, color=TEXT_DARK, bold=False,
                     align=PP_ALIGN.LEFT, space_after=6, level=0):
    p = tf.paragraphs[0] if not tf.paragraphs[0].text and len(tf.paragraphs) == 1 else tf.add_paragraph()
    if tf.paragraphs[0].text == "" and len(tf.paragraphs) == 1:
        p = tf.paragraphs[0]
    p.text = text
    p.level = level
    p.alignment = align
    p.space_after = Pt(space_after)
    if p.runs:
        _set_font(p.runs[0], size, color, bold)
    else:
        _set_font(p.add_run(), size, color, bold)
    return p


def _slide_number(prs, slide, n):
    total = len(prs.slides)
    box = _textbox(slide, SLIDE_W - Inches(1.1), SLIDE_H - Inches(0.42), Inches(0.5), Inches(0.3))
    tf = box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = f"{n}"
    p.alignment = PP_ALIGN.RIGHT
    _set_font(p.runs[0] if p.runs else p.add_run(), 9, TEXT_LIGHT)


def _footer(slide, slide_num=None):
    """Bandeau footer avec nom de l'auteur."""
    bar = _add_rect(slide, 0, SLIDE_H - FOOTER_H, SLIDE_W, FOOTER_H, NAVY)
    # Ligne accent
    _add_rect(slide, 0, SLIDE_H - FOOTER_H, SLIDE_W, Inches(0.04), ACCENT)

    box = _textbox(slide, MARGIN, SLIDE_H - Inches(0.38), Inches(6), Inches(0.3))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{AUTEUR}  ·  {PROJET}  ·  {DATE}"
    _set_font(p.runs[0] if p.runs else p.add_run(), 9, TEXT_LIGHT)


def _header(slide, title, subtitle=None):
    """Bandeau header professionnel."""
    _add_rect(slide, 0, 0, SLIDE_W, HEADER_H, NAVY_MID)
    _add_rect(slide, 0, HEADER_H - Inches(0.05), SLIDE_W, Inches(0.05), ACCENT)

    # Titre
    box = _textbox(slide, MARGIN, Inches(0.18), Inches(9), Inches(0.55))
    tf = box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = title
    _set_font(p.runs[0] if p.runs else p.add_run(), 26, WHITE, bold=True)

    if subtitle:
        box2 = _textbox(slide, MARGIN, Inches(0.68), Inches(9), Inches(0.3))
        tf2 = box2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        _set_font(p2.runs[0] if p2.runs else p.add_run(), 12, ACCENT2)


def _new_content_slide(prs, title, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _fill_slide_bg(slide, LIGHT)
    _header(slide, title, subtitle)
    _footer(slide)
    return slide


def _fill_slide_bg(slide, color):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


# ── Types de slides ─────────────────────────────────────────────────

def add_cover_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_slide_bg(slide, NAVY)

    # Bande décorative gauche
    _add_rect(slide, 0, 0, Inches(0.18), SLIDE_H, ACCENT)

    # Grand rectangle accent en haut à droite
    _add_rect(slide, SLIDE_W - Inches(3.5), 0, Inches(3.5), Inches(0.08), ACCENT2)

    # Titre principal
    box = _textbox(slide, Inches(0.9), Inches(1.6), Inches(10), Inches(1.2))
    tf = box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = PROJET
    _set_font(p.runs[0] if p.runs else p.add_run(), 54, WHITE, bold=True)

    # Sous-titre
    box2 = _textbox(slide, Inches(0.9), Inches(2.7), Inches(10), Inches(0.9))
    tf2 = box2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "Gestion du parc de matériel\nd'un laboratoire technologique"
    _set_font(p2.runs[0] if p2.runs else p2.add_run(), 22, ACCENT2)

    # Ligne séparatrice
    _add_rect(slide, Inches(0.9), Inches(3.75), Inches(2.5), Inches(0.04), ACCENT)

    # Infos auteur
    box3 = _textbox(slide, Inches(0.9), Inches(4.0), Inches(8), Inches(2.2))
    tf3 = box3.text_frame
    tf3.word_wrap = True
    infos = [
        (f"Réalisé par  {AUTEUR}", 18, WHITE, True),
        ("Projet Java — Programmation Orientée Objet", 15, TEXT_LIGHT, False),
        (DATE, 14, TEXT_LIGHT, False),
    ]
    for i, (text, size, color, bold) in enumerate(infos):
        p = tf3.paragraphs[0] if i == 0 else tf3.add_paragraph()
        p.text = text
        p.space_after = Pt(10)
        _set_font(p.runs[0] if p.runs else p.add_run(), size, color, bold)

    # Badge technologie en bas à droite
    badge = _add_rounded(slide, SLIDE_W - Inches(3.2), SLIDE_H - Inches(1.5),
                         Inches(2.6), Inches(0.9), NAVY_MID)
    box4 = _textbox(slide, SLIDE_W - Inches(3.1), SLIDE_H - Inches(1.38),
                    Inches(2.4), Inches(0.7))
    tf4 = box4.text_frame
    p4 = tf4.paragraphs[0]
    p4.text = "Java 17  ·  Swing  ·  SQLite"
    p4.alignment = PP_ALIGN.CENTER
    _set_font(p4.runs[0] if p4.runs else p.add_run(), 11, ACCENT2, bold=True)

    return slide


def add_agenda_slide(prs):
    slide = _new_content_slide(prs, "Plan de la présentation")
    items = [
        ("01", "Contexte & Problématique"),
        ("02", "Objectifs du projet"),
        ("03", "Technologies & Architecture"),
        ("04", "Conception Orientée Objet"),
        ("05", "Architecture MVC"),
        ("06", "Fonctionnalités & Démonstration"),
        ("07", "Structure du code & Persistance"),
        ("08", "Conclusion & Perspectives"),
    ]
    col_w = Inches(5.5)
    for i, (num, label) in enumerate(items):
        col = i % 2
        row = i // 2
        x = MARGIN + col * (col_w + Inches(0.4))
        y = Inches(1.3) + row * Inches(0.72)

        # Numéro
        num_box = _add_rounded(slide, x, y, Inches(0.55), Inches(0.55), ACCENT)
        nb = _textbox(slide, x + Inches(0.05), y + Inches(0.08), Inches(0.45), Inches(0.4))
        tf = nb.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.alignment = PP_ALIGN.CENTER
        _set_font(p.runs[0] if p.runs else p.add_run(), 13, WHITE, bold=True)

        # Label
        lb = _textbox(slide, x + Inches(0.7), y + Inches(0.1), col_w - Inches(0.8), Inches(0.4))
        tf2 = lb.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = label
        _set_font(p2.runs[0] if p2.runs else p.add_run(), 15, TEXT_DARK)

    return slide


def add_bullet_slide(prs, title, bullets, subtitle=None, highlight_idx=None):
    slide = _new_content_slide(prs, title, subtitle)
    box = _textbox(slide, MARGIN, Inches(1.25), SLIDE_W - 2 * MARGIN, SLIDE_H - HEADER_H - FOOTER_H - Inches(0.3))
    tf = box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        is_highlight = highlight_idx and i in highlight_idx
        color = NAVY_MID if is_highlight else TEXT_MID
        bold = is_highlight
        size = 17 if is_highlight else 16

        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"    {bullet}"
        p.space_after = Pt(10)
        p.level = 0

        # Puce colorée via espace + texte
        if p.runs:
            p.runs[0].text = f"●  {bullet}"
        else:
            p.add_run().text = f"●  {bullet}"

        run = p.runs[0]
        _set_font(run, size, color, bold)

        # Colorer la puce
        # (python-pptx ne permet pas de colorer partiellement facilement — on colore tout)
        if not is_highlight:
            _set_font(run, size, TEXT_MID, False)

    return slide


def add_card_slide(prs, title, cards, subtitle=None):
    """Slide avec cartes visuelles (2 ou 3 colonnes)."""
    slide = _new_content_slide(prs, title, subtitle)
    n = len(cards)
    card_w = (SLIDE_W - 2 * MARGIN - Inches(0.3) * (n - 1)) / n
    card_h = SLIDE_H - HEADER_H - FOOTER_H - Inches(0.5)

    for i, (card_title, items, card_color) in enumerate(cards):
        x = MARGIN + i * (card_w + Inches(0.3))
        y = Inches(1.2)

        card = _add_rounded(slide, x, y, card_w, card_h, WHITE)
        card.shadow.inherit = False

        # Bandeau coloré en haut de la carte
        _add_rect(slide, x, y, card_w, Inches(0.55), card_color)

        # Titre carte
        tb = _textbox(slide, x + Inches(0.15), y + Inches(0.1), card_w - Inches(0.3), Inches(0.4))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = card_title
        p.alignment = PP_ALIGN.CENTER
        _set_font(p.runs[0] if p.runs else p.add_run(), 14, WHITE, bold=True)

        # Items
        ib = _textbox(slide, x + Inches(0.2), y + Inches(0.7), card_w - Inches(0.4), card_h - Inches(0.85))
        itf = ib.text_frame
        itf.word_wrap = True
        for j, item in enumerate(items):
            p = itf.paragraphs[0] if j == 0 else itf.add_paragraph()
            p.text = f"▸  {item}"
            p.space_after = Pt(8)
            _set_font(p.runs[0] if p.runs else p.add_run(), 13, TEXT_MID)

    return slide


def add_hierarchy_slide(prs):
    slide = _new_content_slide(prs, "Conception Orientée Objet", "Hiérarchie des classes")

    # Classe abstraite (centre haut)
    cx = SLIDE_W / 2 - Inches(2.5)
    parent = _add_rounded(slide, cx, Inches(1.3), Inches(5), Inches(0.9), NAVY_MID)
    pb = _textbox(slide, cx + Inches(0.1), Inches(1.42), Inches(4.8), Inches(0.7))
    tf = pb.text_frame
    p = tf.paragraphs[0]
    p.text = "«abstract» Ressource"
    p.alignment = PP_ALIGN.CENTER
    _set_font(p.runs[0] if p.runs else p.add_run(), 16, WHITE, bold=True)
    p2 = tf.add_paragraph()
    p2.text = "id · désignation · emplacement · quantité"
    p2.alignment = PP_ALIGN.CENTER
    _set_font(p2.runs[0] if p2.runs else p2.add_run(), 10, ACCENT2)

    # Connecteur
    _add_rect(slide, SLIDE_W / 2 - Inches(0.02), Inches(2.2), Inches(0.04), Inches(0.4), ACCENT)

    # 3 sous-classes
    children = [
        ("EquipementAutomate", "marque · nb E/S · protocole", ACCENT),
        ("EquipementDrone", "masse · autonomie · capteur", RGBColor(0x06, 0x9, 0x8A)),
        ("DocumentationTechnique", "auteur · pages · lien PDF", RGBColor(0x7C, 0x3A, 0xED)),
    ]
    card_w = Inches(3.6)
    start_x = MARGIN
    gap = (SLIDE_W - 2 * MARGIN - 3 * card_w) / 2

    for i, (name, attrs, color) in enumerate(children):
        x = start_x + i * (card_w + gap)
        y = Inches(2.7)

        _add_rounded(slide, x, y, card_w, Inches(1.1), color)
        cb = _textbox(slide, x + Inches(0.1), y + Inches(0.12), card_w - Inches(0.2), Inches(0.9))
        ctf = cb.text_frame
        ctf.word_wrap = True
        cp = ctf.paragraphs[0]
        cp.text = name
        cp.alignment = PP_ALIGN.CENTER
        _set_font(cp.runs[0] if cp.runs else cp.add_run(), 13, WHITE, bold=True)
        cp2 = ctf.add_paragraph()
        cp2.text = attrs
        cp2.alignment = PP_ALIGN.CENTER
        _set_font(cp2.runs[0] if cp2.runs else cp2.add_run(), 10, WHITE)

        # Ligne connecteur
        _add_rect(slide, x + card_w / 2 - Inches(0.02), Inches(2.6), Inches(0.04), Inches(0.12), ACCENT)

    # Encadré méthodes abstraites
    info = _add_rounded(slide, MARGIN, Inches(4.1), SLIDE_W - 2 * MARGIN, Inches(1.5), RGBColor(0xE0, 0xF2, 0xFE))
    ib = _textbox(slide, MARGIN + Inches(0.25), Inches(4.25), SLIDE_W - 2 * MARGIN - Inches(0.5), Inches(1.2))
    itf = ib.text_frame
    itf.word_wrap = True
    ip = itf.paragraphs[0]
    ip.text = "Méthodes abstraites :  getType()  ·  getDiagnostic()"
    _set_font(ip.runs[0] if ip.runs else ip.add_run(), 14, NAVY_MID, bold=True)
    ip2 = itf.add_paragraph()
    ip2.text = "Chaque sous-classe redéfinit ces méthodes → polymorphisme à l'exécution"
    ip2.space_before = Pt(6)
    _set_font(ip2.runs[0] if ip2.runs else ip2.add_run(), 13, TEXT_MID)

    return slide


def add_mvc_slide(prs):
    slide = _new_content_slide(prs, "Architecture MVC", "Séparation des responsabilités")

    layers = [
        ("VUE", "FenetrePrincipale", "Swing · JTable · Formulaire adaptatif", ACCENT),
        ("CONTRÔLEUR", "RessourceController", "Validation · CRUD · Filtrage", NAVY_MID),
        ("MODÈLE", "Ressource + sous-classes", "Données métier · POO", RGBColor(0x06, 0x9, 0x8A)),
        ("DAO", "RessourceDAO · DatabaseManager", "JDBC · SQLite · PreparedStatement", RGBColor(0x7C, 0x3A, 0xED)),
    ]

    box_w = SLIDE_W - 2 * MARGIN - Inches(1.5)
    x = MARGIN + Inches(1.5)
    arrow_x = MARGIN + Inches(0.3)

    for i, (layer, cls, desc, color) in enumerate(layers):
        y = Inches(1.25) + i * Inches(1.15)

        card = _add_rounded(slide, x, y, box_w, Inches(0.95), WHITE)
        _add_rect(slide, x, y, Inches(0.12), Inches(0.95), color)

        # Label couche
        lb = _textbox(slide, x + Inches(0.25), y + Inches(0.08), Inches(2.2), Inches(0.35))
        ltf = lb.text_frame
        lp = ltf.paragraphs[0]
        lp.text = layer
        _set_font(lp.runs[0] if lp.runs else lp.add_run(), 11, color, bold=True)

        # Classe
        cb = _textbox(slide, x + Inches(0.25), y + Inches(0.38), Inches(4), Inches(0.35))
        ctf = cb.text_frame
        cp = ctf.paragraphs[0]
        cp.text = cls
        _set_font(cp.runs[0] if cp.runs else cp.add_run(), 14, TEXT_DARK, bold=True)

        # Description
        db = _textbox(slide, x + Inches(4.5), y + Inches(0.3), box_w - Inches(4.8), Inches(0.5))
        dtf = db.text_frame
        dp = dtf.paragraphs[0]
        dp.text = desc
        _set_font(dp.runs[0] if dp.runs else dp.add_run(), 12, TEXT_MID)

        # Flèche vers le bas
        if i < len(layers) - 1:
            ab = _textbox(slide, arrow_x, y + Inches(0.95), Inches(0.8), Inches(0.25))
            atf = ab.text_frame
            ap = atf.paragraphs[0]
            ap.text = "▼"
            ap.alignment = PP_ALIGN.CENTER
            _set_font(ap.runs[0] if ap.runs else ap.add_run(), 14, ACCENT, bold=True)

    # Flux
    flux = _textbox(slide, MARGIN, SLIDE_H - FOOTER_H - Inches(0.55), SLIDE_W - 2 * MARGIN, Inches(0.4))
    ftf = flux.text_frame
    fp = ftf.paragraphs[0]
    fp.text = "Flux :  Utilisateur  →  Vue  →  Contrôleur  →  DAO  →  Base SQLite"
    fp.alignment = PP_ALIGN.CENTER
    _set_font(fp.runs[0] if fp.runs else fp.add_run(), 12, NAVY_MID, bold=True)

    return slide


def add_features_slide(prs):
    slide = _new_content_slide(prs, "Fonctionnalités clés")

    features = [
        ("📋", "Tableau dynamique", "Affichage complet du stock avec diagnostic polymorphe"),
        ("📝", "Formulaire adaptatif", "Champs activés/désactivés selon le type de ressource"),
        ("🔍", "Recherche temps réel", "Filtrage instantané par désignation, type ou emplacement"),
        ("✅", "Validation robuste", "Messages d'erreur explicites via JOptionPane"),
        ("💾", "Persistance SQLite", "Sauvegarde automatique en base locale JDBC"),
        ("🗑️", "Suppression sécurisée", "Confirmation avant suppression avec dialogue"),
    ]

    cols = 2
    card_w = (SLIDE_W - 2 * MARGIN - Inches(0.3)) / cols
    card_h = Inches(1.35)

    for i, (icon, feat_title, feat_desc) in enumerate(features):
        col = i % cols
        row = i // cols
        x = MARGIN + col * (card_w + Inches(0.3))
        y = Inches(1.25) + row * (card_h + Inches(0.2))

        _add_rounded(slide, x, y, card_w, card_h, WHITE)

        # Icône
        ib = _textbox(slide, x + Inches(0.15), y + Inches(0.2), Inches(0.5), Inches(0.5))
        itf = ib.text_frame
        ip = itf.paragraphs[0]
        ip.text = icon
        _set_font(ip.runs[0] if ip.runs else ip.add_run(), 22, ACCENT)

        # Titre
        tb = _textbox(slide, x + Inches(0.7), y + Inches(0.15), card_w - Inches(0.85), Inches(0.4))
        ttf = tb.text_frame
        tp = ttf.paragraphs[0]
        tp.text = feat_title
        _set_font(tp.runs[0] if tp.runs else tp.add_run(), 14, TEXT_DARK, bold=True)

        # Description
        db = _textbox(slide, x + Inches(0.7), y + Inches(0.55), card_w - Inches(0.85), Inches(0.7))
        dtf = db.text_frame
        dtf.word_wrap = True
        dp = dtf.paragraphs[0]
        dp.text = feat_desc
        _set_font(dp.runs[0] if dp.runs else dp.add_run(), 11, TEXT_MID)

    return slide


def add_ui_slide(prs):
    slide = _new_content_slide(prs, "Interface utilisateur", "3 zones fonctionnelles")

    zones = [
        ("Zone 1 — Recherche", "Champ de filtrage en temps réel\nsur type, désignation, emplacement", ACCENT),
        ("Zone 2 — Tableau", "JTable avec tri et colonnes :\nID · Type · Désignation · Emplacement · Qté · Diagnostic", NAVY_MID),
        ("Zone 3 — Formulaire", "Ajout de ressource avec champs\nadaptatifs selon le type sélectionné", RGBColor(0x06, 0x9, 0x8A)),
    ]

    for i, (zone_title, zone_desc, color) in enumerate(zones):
        y = Inches(1.25) + i * Inches(1.55)
        _add_rounded(slide, MARGIN, y, SLIDE_W - 2 * MARGIN, Inches(1.35), WHITE)
        _add_rect(slide, MARGIN, y, Inches(0.12), Inches(1.35), color)

        tb = _textbox(slide, MARGIN + Inches(0.3), y + Inches(0.12), Inches(4), Inches(0.4))
        ttf = tb.text_frame
        tp = ttf.paragraphs[0]
        tp.text = zone_title
        _set_font(tp.runs[0] if tp.runs else tp.add_run(), 15, color, bold=True)

        db = _textbox(slide, MARGIN + Inches(0.3), y + Inches(0.52), Inches(7), Inches(0.75))
        dtf = db.text_frame
        dtf.word_wrap = True
        for j, line in enumerate(zone_desc.split("\n")):
            dp = dtf.paragraphs[0] if j == 0 else dtf.add_paragraph()
            dp.text = line
            _set_font(dp.runs[0] if dp.runs else dp.add_run(), 12, TEXT_MID)

    # Placeholder capture
    ph = _add_rounded(slide, SLIDE_W - MARGIN - Inches(3.8), Inches(1.25),
                      Inches(3.5), Inches(4.6), RGBColor(0xE2, 0xE8, 0xF0))
    pb = _textbox(slide, SLIDE_W - MARGIN - Inches(3.65), Inches(3.0), Inches(3.2), Inches(1.2))
    ptf = pb.text_frame
    ptf.word_wrap = True
    pp = ptf.paragraphs[0]
    pp.text = "[ Capture d'écran ]"
    pp.alignment = PP_ALIGN.CENTER
    _set_font(pp.runs[0] if pp.runs else pp.add_run(), 13, TEXT_LIGHT, italic=True)
    pp2 = ptf.add_paragraph()
    pp2.text = "Insérer ici une\nimage de l'application"
    pp2.alignment = PP_ALIGN.CENTER
    _set_font(pp2.runs[0] if pp2.runs else pp2.add_run(), 11, TEXT_LIGHT, italic=True)

    return slide


def add_code_structure_slide(prs):
    slide = _new_content_slide(prs, "Structure du code source", "10 fichiers Java · ~1 200 lignes")

    packages = [
        ("com.glab", "GLabApp.java", "Point d'entrée", NAVY),
        ("com.glab.model", "4 classes POO", "Ressource (abstract) + 3 sous-classes", ACCENT),
        ("com.glab.controller", "RessourceController", "Validation · Ajout · Suppression · Filtrage", NAVY_MID),
        ("com.glab.database", "DatabaseManager · RessourceDAO", "Connexion SQLite · CRUD polymorphe", RGBColor(0x06, 0x9, 0x8A)),
        ("com.glab.view", "FenetrePrincipale", "Interface Swing complète", RGBColor(0x7C, 0x3A, 0xED)),
    ]

    for i, (pkg, files, role, color) in enumerate(packages):
        y = Inches(1.2) + i * Inches(0.95)
        _add_rounded(slide, MARGIN, y, SLIDE_W - 2 * MARGIN, Inches(0.82), WHITE)
        _add_rect(slide, MARGIN, y, Inches(0.1), Inches(0.82), color)

        pb = _textbox(slide, MARGIN + Inches(0.2), y + Inches(0.08), Inches(3.5), Inches(0.35))
        ptf = pb.text_frame
        pp = ptf.paragraphs[0]
        pp.text = pkg
        _set_font(pp.runs[0] if pp.runs else pp.add_run(), 12, color, bold=True)

        fb = _textbox(slide, MARGIN + Inches(3.8), y + Inches(0.08), Inches(3.5), Inches(0.35))
        ftf = fb.text_frame
        fp = ftf.paragraphs[0]
        fp.text = files
        _set_font(fp.runs[0] if fp.runs else fp.add_run(), 12, TEXT_DARK, bold=True)

        rb = _textbox(slide, MARGIN + Inches(0.2), y + Inches(0.42), SLIDE_W - 2 * MARGIN - Inches(0.4), Inches(0.35))
        rtf = rb.text_frame
        rp = rtf.paragraphs[0]
        rp.text = role
        _set_font(rp.runs[0] if rp.runs else rp.add_run(), 11, TEXT_MID)

    # Commandes
    cmd_box = _add_rounded(slide, MARGIN, Inches(6.0), SLIDE_W - 2 * MARGIN, Inches(0.65),
                           RGBColor(0x1E, 0x29, 0x3B))
    cb = _textbox(slide, MARGIN + Inches(0.2), Inches(6.1), SLIDE_W - 2 * MARGIN - Inches(0.4), Inches(0.45))
    ctf = cb.text_frame
    cp = ctf.paragraphs[0]
    cp.text = "$ mvn exec:java   ·   bash run.sh   ·   java -jar target/g-lab-1.0.0.jar"
    _set_font(cp.runs[0] if cp.runs else cp.add_run(), 12, ACCENT2, bold=True)

    return slide


def add_conclusion_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_slide_bg(slide, NAVY)
    _add_rect(slide, 0, 0, Inches(0.18), SLIDE_H, ACCENT)

    # Titre
    box = _textbox(slide, Inches(0.9), Inches(1.2), Inches(10), Inches(0.9))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = "Conclusion"
    _set_font(p.runs[0] if p.runs else p.add_run(), 40, WHITE, bold=True)

    _add_rect(slide, Inches(0.9), Inches(2.1), Inches(2.5), Inches(0.04), ACCENT)

    # Points clés
    points = [
        "Application concrète des 4 piliers de la POO",
        "Architecture MVC maintenable et évolutive",
        "Persistance fiable via JDBC / SQLite",
        "Interface ergonomique avec formulaire adaptatif",
    ]
    pb = _textbox(slide, Inches(0.9), Inches(2.4), Inches(7), Inches(2.5))
    ptf = pb.text_frame
    ptf.word_wrap = True
    for i, pt in enumerate(points):
        pp = ptf.paragraphs[0] if i == 0 else ptf.add_paragraph()
        pp.text = f"✔   {pt}"
        pp.space_after = Pt(12)
        _set_font(pp.runs[0] if pp.runs else pp.add_run(), 16, ACCENT2)

    # Perspectives
    pers_box = _add_rounded(slide, Inches(0.9), Inches(4.9), Inches(7), Inches(1.2), NAVY_MID)
    pers_b = _textbox(slide, Inches(1.1), Inches(5.0), Inches(6.6), Inches(1.0))
    pers_tf = pers_b.text_frame
    pers_tf.word_wrap = True
    pp1 = pers_tf.paragraphs[0]
    pp1.text = "Perspectives"
    _set_font(pp1.runs[0] if pp1.runs else pp1.add_run(), 13, GOLD, bold=True)
    pp2 = pers_tf.add_paragraph()
    pp2.text = "Nouveaux types de ressources · Migration PostgreSQL · Interface web · Export PDF"
    _set_font(pp2.runs[0] if pp2.runs else pp2.add_run(), 12, TEXT_LIGHT)

    # Merci + nom
    mb = _textbox(slide, Inches(0.9), Inches(6.2), Inches(10), Inches(0.9))
    mtf = mb.text_frame
    mp = mtf.paragraphs[0]
    mp.text = "Merci pour votre attention"
    _set_font(mp.runs[0] if mp.runs else mp.add_run(), 22, WHITE, bold=True)
    mp2 = mtf.add_paragraph()
    mp2.text = AUTEUR
    mp2.space_before = Pt(4)
    _set_font(mp2.runs[0] if mp2.runs else mp2.add_run(), 16, ACCENT2, bold=True)
    mp3 = mtf.add_paragraph()
    mp3.text = "Questions ?"
    mp3.space_before = Pt(6)
    _set_font(mp3.runs[0] if mp3.runs else mp3.add_run(), 14, TEXT_LIGHT, italic=True)

    return slide


# ── Main ─────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    add_cover_slide(prs)
    add_agenda_slide(prs)

    add_bullet_slide(prs, "Contexte & Problématique", [
        "Laboratoire universitaire avec un parc hétérogène de matériel technologique",
        "Automates industriels, drones embarqués et documentation technique à gérer",
        "La gestion manuelle du stock entraîne des erreurs et une perte de temps",
        "Besoin d'une application centralisée, persistante et ergonomique",
        "Solution développée : G-Lab — application Java desktop avec interface Swing",
    ])

    add_bullet_slide(prs, "Objectifs du projet", [
        "Modéliser les ressources par une hiérarchie de classes orientée objet",
        "Appliquer les 4 piliers POO : abstraction, encapsulation, héritage, polymorphisme",
        "Implémenter une architecture MVC claire et maintenable",
        "Valider les saisies utilisateur avec des messages d'erreur explicites",
        "Assurer la persistance des données via JDBC et SQLite",
    ])

    add_card_slide(prs, "Technologies & Outils", [
        ("Backend", ["Java 17 (LTS)", "JDBC + SQLite 3.45", "Maven — build & dépendances", "SLF4J — journalisation"], NAVY_MID),
        ("Interface", ["Swing — JTable, JFrame", "JOptionPane — dialogues", "Formulaire adaptatif", "Filtrage temps réel"], ACCENT),
        ("Architecture", ["Pattern MVC", "Pattern DAO", "Try-with-resources", "PreparedStatement"], RGBColor(0x06, 0x9, 0x8A)),
    ])

    add_hierarchy_slide(prs)

    add_card_slide(prs, "Les 4 piliers de la POO", [
        ("Abstraction", ["Classe Ressource abstraite", "Contrat via getType() et getDiagnostic()", "Instanciation impossible directement"], NAVY_MID),
        ("Encapsulation", ["Attributs private", "Accesseurs / mutateurs publics", "Validation dans le contrôleur"], ACCENT),
        ("Héritage", ["3 sous-classes de Ressource", "Réutilisation via super()", "Spécialisation des attributs"], RGBColor(0x06, 0x9, 0x8A)),
        ("Polymorphisme", ["getDiagnostic() redéfini", "DAO accepte tout type Ressource", "Switch de reconstruction en base"], RGBColor(0x7C, 0x3A, 0xED)),
    ], subtitle="Principes fondamentaux appliqués dans G-Lab")

    add_mvc_slide(prs)
    add_features_slide(prs)
    add_ui_slide(prs)
    add_code_structure_slide(prs)

    add_bullet_slide(prs, "Persistance des données", [
        "Base SQLite locale glab_parc.db — créée automatiquement au premier lancement",
        "Table unique ressources avec colonnes génériques pour les attributs spécifiques",
        "PreparedStatement pour se protéger contre les injections SQL",
        "Try-with-resources pour la fermeture automatique des connexions JDBC",
        "Pattern DAO : RessourceDAO isole tout accès à la base de données",
        "Reconstruction polymorphe des objets à la lecture via switch sur le type",
    ], subtitle="JDBC · SQLite · Pattern DAO")

    add_conclusion_slide(prs)

    prs.save(OUTPUT)
    print(f"✔ Présentation professionnelle créée : {OUTPUT}")
    print(f"  Auteur : {AUTEUR}")
    print(f"  Diapositives : {len(prs.slides)}")


if __name__ == "__main__":
    main()
